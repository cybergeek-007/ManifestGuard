from pptx import Presentation

def replace_text(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            new_text = original_text
            for old, new in replacements.items():
                new_text = new_text.replace(old, new)
            if new_text != original_text:
                run.text = new_text

def process_ppt(filepath):
    prs = Presentation(filepath)
    
    global_replacements = {
        "POWER SCORE": "REACH SCORE",
        "SUSPICION SCORE": "ANOMALY SCORE",
        "POWER": "REACH",
        "SUSPICION": "ANOMALY",
        "Power": "Reach",
        "Suspicion": "Anomaly",
        "power_score": "reach_score",
        "suspicion_score": "anomaly_score",
        "Bitwarden": "uBlock Origin",
        "Password Manager": "Ad Blocker"
    }
    
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text(shape, global_replacements)
                
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text(cell, global_replacements)

    # Specific slide updates
    
    # Slide 3 (Index 2): Update the 4th dimension to Threat Intel & Collusion instead of Allowlist
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "COMMUNITY ALLOWLIST" in run.text:
                        run.text = run.text.replace("COMMUNITY ALLOWLIST", "SUPPLY CHAIN & INTEL")
                    if "Is it a verified, known-good tool?" in run.text:
                        run.text = run.text.replace("Is it a verified, known-good tool?", "Are versions safe & domains clean?")
                    if "200+ curated extensions" in run.text:
                        run.text = run.text.replace("200+ curated extensions ? safe alternative recommendations", "Version delta tracking \u2022 domain intel burst \u2022 collusion graph")

    # Slide 6 (Index 5): Update verdict ladder table
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_table:
            table = shape.table
            # Let's clear and refill the table or just replace cell text
            # V4 ladder: trusted, low_concern, moderate_risk, suspicious, known_malicious
            # Current rows:
            # Header
            # known_malicious
            # disabled_by_chrome (to be replaced with moderate_risk)
            # suspicious
            # removed_or_unavailable (to be removed or replaced with low_concern)
            # powerful_but_expected (removed)
            # trusted
            
            # Since table row deletion in python-pptx is tricky, we can just overwrite the 5 rows below header.
            # And blank out any extra rows.
            v4_ladder = [
                ("Verdict", "Meaning", "Trigger"),
                ("known_malicious", "Matched threat intelligence DB", "Intel match or CWS unavail"),
                ("suspicious", "Abnormal code behavior detected", "Anomaly \u2265 40"),
                ("moderate_risk", "High reach without strong reputation", "Reach \u2265 40 & Rep < 40"),
                ("low_concern", "Generally safe with minor signals", "Anomaly < 40 & Reach < 40"),
                ("trusted", "High reputation or allowlisted", "Allowlist or Rep \u2265 70")
            ]
            
            for row_idx, row_data in enumerate(v4_ladder):
                if row_idx < len(table.rows):
                    for col_idx, text in enumerate(row_data):
                        table.cell(row_idx, col_idx).text = text
                        
            # If there are more rows, blank them
            for row_idx in range(len(v4_ladder), len(table.rows)):
                for cell in table.rows[row_idx].cells:
                    cell.text = ""

    prs.save("D:/ManifestGuard/ManifestGuard_v4.pptx")
    print("Saved as ManifestGuard_v4.pptx")

process_ppt("D:/ManifestGuard/ManifestGuard.pptx")
