from pptx import Presentation

def dump_ppt_text(filepath):
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(run.text.encode('ascii', 'replace').decode('ascii'))

dump_ppt_text("D:/ManifestGuard/ManifestGuard.pptx")
